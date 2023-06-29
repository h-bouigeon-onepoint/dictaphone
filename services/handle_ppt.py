from pptx import Presentation
from services.constants import Constants
from domain.model.meeting_report_data import MeetingReportData


class PPTWriter:
    def __init__(self, filename: str, meeting_report: MeetingReportData):
        self.filename = filename
        self.meeting_report = meeting_report
        self.ppt_path = None

    def write_to_ppt(self, ppt_path: str) -> None:
        # Read template
        presentation = Presentation(Constants.TEMPLATE_PPT_PATH)
        slide_to_modify = presentation.slides[0]

        # Fill elements in slide
        if self.meeting_report != Constants.MODEL_JSON_MEETING_REPORT:
            self._fill_title_meeting(slide_to_modify)
            self._fill_participants(slide_to_modify)
            self._fill_agenda(slide_to_modify)
            self._fill_elements_discussed(slide_to_modify)
            self._fill_notes(slide_to_modify)
            self._fill_actions(slide_to_modify, presentation)

        presentation.save(ppt_path)

    def _fill_title_meeting(self, slide_to_modify) -> None:
        title_meeting_cell = slide_to_modify.shapes[7]
        title_meeting_cell.text = self.meeting_report.title_meeting

    def _fill_participants(self, slide_to_modify) -> None:
        participants_cell = slide_to_modify.shapes[6].table.cell(1, 0)
        participants_cell.text = "- " + "\n- ".join(self.meeting_report.participants)

    def _fill_agenda(self, slide_to_modify) -> None:
        agenda_cell = slide_to_modify.shapes[3].table.cell(1, 0)
        agenda_cell.text = "• " + "\n• ".join(self.meeting_report.agenda)

    def _fill_elements_discussed(self, slide_to_modify) -> None:
        elements_discussed_cell = slide_to_modify.shapes[4].table.cell(1, 0)
        elements_discussed_cell.text = "• " + "\n• ".join(
            self.meeting_report.elements_discussed[
                :7
            ]  # limit 7 elements max for nicer display
        )

    def _fill_notes(self, slide_to_modify) -> None:
        notes_cell = slide_to_modify.shapes[8].table.cell(1, 0)
        notes_cell.text = self.meeting_report.notes

    def _fill_actions(self, slide_to_modify, presentation) -> None:
        actions_table = presentation.slides[0].shapes[5].table
        actions = self.meeting_report.action_to_be_taken

        if actions != []:
            for i in range(len(actions)):
                if (
                    i - 1 > len(actions_table.rows) or i > 4
                ):  # limit 4 elements max for nicer display
                    break
                actions_cell = slide_to_modify.shapes[5].table.cell(i + 1, 0)
                actions_cell.text = actions[i]["action"]

                actions_manager_cell = (
                    presentation.slides[0].shapes[5].table.cell(i + 1, 1)
                )
                actions_manager_cell.text = actions[i]["responsable"]

                actions_due_date_cell = (
                    presentation.slides[0].shapes[5].table.cell(i + 1, 2)
                )
                actions_due_date_cell.text = actions[i]["deadline"]
